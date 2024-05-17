from argparse import ArgumentParser
parser = ArgumentParser()
parser.add_argument("--pptx", type=str, required=True)
parser.add_argument("--output", type=str,required=True)
parser.add_argument("--slide", type=int,default=0)
args = parser.parse_args()

from pptx import Presentation
import pptx
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image
from io import BytesIO
from tqdm import tqdm

import json

path = args.pptx
prs = Presentation(path)
patches = list()
emu2cm = 1 / 360000
emu2inch = 1 / 914400
slide = prs.slides[args.slide]

for i, shape in enumerate(tqdm(list(slide.shapes))):
    if isinstance(shape, pptx.shapes.picture.Picture):
        try:
            img = np.flipud(Image.open(BytesIO(shape.image.blob)))
            extent = shape.left, shape.top, shape.width, shape.height   
            extent = [emu2inch*e for e in extent]
            print("DBG", extent)    
            patches.append((extent, shape.image.dpi, np.array(img), shape.rotation))
        except:
            print("Load img error")

from copy import deepcopy
save = dict()
save["patches"] = []
save["background"] = deepcopy({"x": 0, "y": 0, "w":8, "h":6, "dpi": 150, "rotation": 0})
save["optical"] = deepcopy({"x": 0, "y": 0, "w":8, "h":6, "dpi": 150, "rotation": 0})
save["global"] = deepcopy({"angle": 0, "scale": 1.0, "folder": "pngs_22fep01_p"})


for i, patch in enumerate(tqdm(patches)):
    extent, dpi, img,rot = patch
    x, y, w, h = extent
    save["patches"].append(deepcopy({"x": x, "y": y, "w":h, "h":w, "dpi": dpi, "rotation": rot}))
    pimg = Image.fromarray(np.fliplr(np.rot90(img)))
    pimg.save(f"{i}.png")

with open(args.output, "w") as f:
    json.dump(save, f)

